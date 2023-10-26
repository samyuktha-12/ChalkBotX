import chainlit as cl
from io import BytesIO
from pptx import Presentation
import chunks

async def ppt_processing(ppt_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload PPT files",
            accept=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
            max_size_mb=20,
            timeout=180,
            max_files=ppt_count
            ).send()
    
    i=0
    for file in files:
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        ppt_stream = BytesIO(file.content)
        presentation = Presentation(ppt_stream)
        ppt_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + "\n"

        # Split the text into chunks
        chunks.chunk_splitter(ppt_text,texts,metadatas,file,i)
        i+=1
    return texts,metadatas