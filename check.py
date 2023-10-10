import os
import io
import chainlit as cl
import PyPDF2
from io import BytesIO
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain.chains import RetrievalQAWithSourcesChain
from langchain.chat_models import ChatOpenAI
from chainlit.input_widget import TextInput
from langchain.prompts.chat import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)

from dotenv import load_dotenv
from pptx import Presentation
import parameters
import scanned_to_searchable
import pytesseract
from PIL import Image
import requests

text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

system_template = """Use the following pieces of context to answer the users question.
If you don't know the answer, just say that you don't know, don't try to make up an answer.
ALWAYS return a "SOURCES" part in your answer.
The "SOURCES" part should be a reference to the source of the document from which you got your answer.

Example of your response should be:

```
The answer is foo
SOURCES: xyz
```

Begin!
----------------
{summaries}"""

messages = [
    SystemMessagePromptTemplate.from_template(system_template),
    HumanMessagePromptTemplate.from_template("{question}"),
]
prompt = ChatPromptTemplate.from_messages(messages)
chain_type_kwargs = {"prompt": prompt}


@cl.on_chat_start
async def start():
    settings = await cl.ChatSettings(
        [
            TextInput(id="pdf", label="Enter number of pdf files", initial="0"),
            TextInput(id="scanned_pdf", label="Enter number of scanned pdf files", initial="0"),
            TextInput(id="txt", label="Enter number of text files", initial="0"),
            TextInput(id="ppt", label="Enter number of PPT files", initial="0"),
            TextInput(id="img", label="Enter number of Image files", initial="0"),
            TextInput(id="url", label="Enter number of URLs", initial="0")

        ]
    ).send()


@cl.on_settings_update
async def setup_agent(settings):
    print("on_settings_update", settings)
    pdf_count=int(settings["pdf"])
    txt_count=int(settings["txt"])
    scanned_pdf_count=int(settings["scanned_pdf"])
    ppt_count=int(settings["ppt"])
    img_count=int(settings["img"])
    url_count=int(settings["url"])

    elements = [
    cl.Image(name="image1", display="inline", path="./robot.jpeg")
    ]
    await cl.Message(content="Hello there, Input all relevent docs!", elements=elements).send()
    files = None

    texts=[]
    metadatas=[]
    if pdf_count!=0:
        await pdf_processing(pdf_count,texts,metadatas)

    if txt_count!=0:
        await txt_processing(txt_count,texts,metadatas)
    
    if scanned_pdf_count!=0:
        await scanned_processing(scanned_pdf_count,texts,metadatas)
    
    if ppt_count!=0:
        await ppt_processing(ppt_count,texts,metadatas)
    
    if img_count!=0:
        await img_processing(img_count,texts,metadatas)

    if url_count!=0:
        await url_processing(url_count,texts,metadatas)

    print(texts)
    # Create a Chroma vector store
    embeddings = OpenAIEmbeddings()
    docsearch = await cl.make_async(Chroma.from_texts)(
        texts, embeddings, metadatas=metadatas
    )

    # Create a chain that uses the Chroma vector store
    chain = RetrievalQAWithSourcesChain.from_chain_type(
        ChatOpenAI(temperature=0),
        chain_type="stuff",
        retriever=docsearch.as_retriever(),
    )

    # Save the metadata and texts in the user session
    cl.user_session.set("metadatas", metadatas)
    cl.user_session.set("texts", texts)
    print(texts)
    print(metadatas)

    # Let the user know that the system is ready
    msg= cl.Message(content=f"Processing done. You can now ask questions!")
    await msg.send()

    cl.user_session.set("chain", chain)



async def pdf_processing(pdf_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload PDF files",
            accept=["application/pdf"],
            max_size_mb=20,
            timeout=180,
            max_files=pdf_count
            ).send()
    
    
    for file in files:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        # Read the PDF file
        pdf_stream = BytesIO(file.content)
        pdf = PyPDF2.PdfReader(pdf_stream)
        pdf_text = ""
        for page in pdf.pages:
            pdf_text += page.extract_text()

        # Split the text into chunks
        chunks=text_splitter.split_text(pdf_text)
        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{file.name}:chunk {i}"})
            i=i+1

async def txt_processing(txt_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload text files",
            accept=["text/plain"],
            max_size_mb=20,
            timeout=180,
            max_files=txt_count
            ).send()
    
    
    for file in files:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        # Decode the file
        text = file.content.decode("utf-8")

        # Split the text into chunks
        chunks = text_splitter.split_text(text)

        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{file.name}:chunk {i}"})
            i=i+1

async def scanned_processing(scanned_pdf_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload Scanned PDF files",
            accept=["application/pdf"],
            max_size_mb=20,
            timeout=180,
            max_files=scanned_pdf_count
            ).send()
    
    for file in files:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        # Read the PDF file
        pdf_stream = BytesIO(file.content)
        pdf = PyPDF2.PdfReader(pdf_stream)
        output_pdf=PyPDF2.PdfWriter()
        for page in pdf.pages:
            output_pdf.add_page(page)
        completeName = os.path.join(parameters.inpath, file.name)
        output_pdf.write(completeName)
        scanned_to_searchable.image_conversion(completeName)
        scanned_to_searchable.pdf_conversion(parameters.folder_dir)
        scanned_to_searchable.merge_pdf(parameters.pdf_folder,file.name)
        for p in parameters.directory_path:
            scanned_to_searchable.delete_files_in_directory(p)

        # Read the Searchable PDF file
        pdf = PyPDF2.PdfReader(parameters.pdfpath+'/'+file.name+".pdf")
        pdf_text = ""
        for page in pdf.pages:
            pdf_text += page.extract_text()

        # Split the text into chunks
        chunks=text_splitter.split_text(pdf_text)
        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{file.name}:chunk {i}"})
            i=i+1
    scanned_to_searchable.delete_files_in_directory(parameters.pdfpath)

async def ppt_processing(ppt_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload PPT files",
            accept=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
            max_size_mb=20,
            timeout=180,
            max_files=ppt_count
            ).send()
    
    
    for file in files:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        ppt_stream = BytesIO(file.content)
        presentation = Presentation(ppt_stream)
        ppt_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + "\n"

        # Split the text into chunks
        chunks=text_splitter.split_text(ppt_text)
        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{file.name}:chunk {i}"})
            i=i+1

async def img_processing(img_count,texts,metadatas):
    files = await cl.AskFileMessage(
            content="Please upload JPG files",
            accept=["image/jpeg"],
            max_size_mb=20,
            timeout=180,
            max_files=img_count
            ).send()
    
    for file in files:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing `{file.name}`...")
        await msg.send()

        # Read the PDF file
        image1=io.BytesIO(file.content)
        image = Image.open(image1)
        extracted_text = pytesseract.image_to_string(image)

        # Split the text into chunks
        chunks=text_splitter.split_text(extracted_text)
        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{file.name}:chunk {i}"})
            i=i+1

async def url_processing(url_count,texts,metadatas):
    urls = await cl.AskUserMessage(
            content="Please upload urls (as comma separated text)",
            ).send()
    urls=urls['content'].split(",")
    for url in urls:
        i=0 #for chunk metadata
        msg = cl.Message(content=f"Processing ...")
        await msg.send()
        
        response = requests.get(url)
        if response.status_code != 200:
            return None, f"Failed to fetch content from the URL: {url}"
        content = response.text
        chunks = content.split("\n\n")  # Split by double line breaks (paragraphs)

        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({"source": f"{url}:chunk {i}"})
            i=i+1

@cl.on_message
async def main(message:str):

    chain = cl.user_session.get("chain")  # type: RetrievalQAWithSourcesChain
    cb = cl.AsyncLangchainCallbackHandler(
        stream_final_answer=True, answer_prefix_tokens=["FINAL", "ANSWER"]
    )
    cb.answer_reached = True
    res = await chain.acall(message, callbacks=[cb])

    answer = res["answer"]
    sources = res["sources"].strip()
    source_elements = []

    # Get the metadata and texts from the user session
    metadatas = cl.user_session.get("metadatas")
    all_sources = [m["source"] for m in metadatas]
    texts = cl.user_session.get("texts")

    if sources:
        found_sources = []

        # Add the sources to the message
        for source in sources.split(","):
            source_name = source.strip()#.replace(".", "")
            # Get the index of the source
            try:
                index = all_sources.index(source_name)
            except ValueError:
                print("value error in finding source")
                continue
            text = texts[index]
            found_sources.append(source_name)
            # Create the text element referenced in the message
            source_elements.append(cl.Text(content=text, name=source_name))

        if found_sources:
            answer += f"\nSources: {', '.join(found_sources)}"
        else:
            answer += "\nNo sources found"

    if cb.has_streamed_final_answer:
        cb.final_stream.elements = source_elements
        await cb.final_stream.update()
    else:
        await cl.Message(content=answer, elements=source_elements).send()
    